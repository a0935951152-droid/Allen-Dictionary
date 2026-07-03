"""知識側建置：sources/ → Obsidian vault（心智圖）+ CAG 詞表（§3，純淨資料側）。

把 `data/knowledge/sources/` 的乾淨來源自動萃成標準知識庫，全地端、零手寫：
- **xlsx 對照表**（錯字→正確字字典）＝最高價值：正解 term + ASR 變體 + 信心 + 類型 + 上下文。
- **docx/pptx/pdf 乾淨文件** → 兩階段 LLM 抽詞+語意（Pass1 心智圖、Pass2 CAG）。
- 輸出：`vault/terms/<詞>.md`（frontmatter + [[同音/變體/領域]] 連結 → Obsidian graph 即心智圖）、
  `vault/domains/<類>.md`、`cag/<domain>.md`（高信任詞表＝判官 CAG 前綴）。

鐵律：知識側**只收 authored 乾淨來源**；逐字稿/字幕/ASR 校正（.srt/.txt 或檔名含字幕/逐字/校正/三語對齊）
含 ASR 錯字會污染知識側，一律排除（_CLEAN_EXT + _DIRTY_NAME 兩道守門）。
"""
from __future__ import annotations

import asyncio
import hashlib
import json
import os
import re

from .. import clients
from ..config import settings
from . import langtag
from .ipa import to_ipa
from .variants import flatten as _flat


def _names_map(e: dict) -> dict:
    """多語對照 map（ISO-639-3 key，DaMuEL 模式，可擴充）。
    既有 names 優先；否則由 en/ja 遷移（不重跑 LLM）。空值不收。"""
    m = dict(e.get("names") or {})
    if e.get("en") and "eng" not in m:
        m["eng"] = e["en"]
    if e.get("ja") and "jpn" not in m:
        m["jpn"] = e["ja"]
    return {k: v for k, v in m.items() if v}

_SRC = os.path.join(settings.data_dir, "knowledge", "sources")
_VAULT = os.path.join(settings.data_dir, "knowledge", "vault")
_CAG = os.path.join(settings.data_dir, "knowledge", "cag")
_MANIFEST = os.path.join(settings.data_dir, "knowledge", ".manifest.json")
_DECISIONS = os.path.join(settings.data_dir, "knowledge", ".decisions.json")   # 每檔人工抽取/排除覆寫（rel→extract|exclude）
_ACOUSTIC = os.path.join(settings.data_dir, "knowledge", "acoustic.json")       # 聲學種子：正解+ASR變體 → 餵 §4.1 音韻索引（§5.6 回灌）

# 人工編輯後可鎖住、不被 build 覆寫的語意欄位（§11 可持續擴充：鎖定欄位不覆寫）。
_LOCKED_FIELDS = ("term", "category", "gloss", "variants", "reason", "names")


# ── 增量 manifest（內容 hash → 每檔抽取快取；只重抽新/改動檔）──────
def _file_hash(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for blk in iter(lambda: f.read(65536), b""):
            h.update(blk)
    return h.hexdigest()


def _load_manifest() -> dict:
    try:
        with open(_MANIFEST, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _save_manifest(m: dict) -> None:
    tmp = _MANIFEST + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(m, f, ensure_ascii=False)
    os.replace(tmp, _MANIFEST)


# ── 每檔抽取決策覆寫（rel → "extract"|"exclude"）：使用者手動凌駕自動判斷 ──────
def _load_decisions() -> dict:
    try:
        with open(_DECISIONS, encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _save_decisions(d: dict) -> None:
    tmp = _DECISIONS + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(d, f, ensure_ascii=False)
    os.replace(tmp, _DECISIONS)


def set_decision(rel: str, decision: str) -> dict:
    """設/清單檔抽取決策。decision: extract|exclude|auto（auto＝清除覆寫、回到自動判斷）。"""
    if decision not in ("extract", "exclude", "auto"):
        raise ValueError("decision 須為 extract / exclude / auto")
    d = _load_decisions()
    if decision == "auto":
        d.pop(rel, None)
    else:
        d[rel] = decision
    _save_decisions(d)
    return {"rel": rel, "decision": decision, "overrides": len(d)}


def move_source(rel: str, action: str) -> dict:
    """搬移來源檔＝排除/抽取（物理位置即決策）：
    - action=exclude → 搬進 sources/_excluded/（build 不抽）
    - action=extract → 搬回 sources 根目錄（build 會抽）
    回 {rel(新路徑), from, action} 或 {error}。"""
    if action not in ("exclude", "extract"):
        raise ValueError("action 須為 exclude / extract")
    src = os.path.normpath(os.path.join(_SRC, rel))
    if not (src + os.sep).startswith(_SRC + os.sep) or not os.path.isfile(src):
        return {"error": f"找不到來源：{rel}"}
    name = os.path.basename(src)
    dst_dir = os.path.join(_SRC, "_excluded") if action == "exclude" else _SRC
    os.makedirs(dst_dir, exist_ok=True)
    dst = os.path.join(dst_dir, name)
    if os.path.normpath(dst) == os.path.normpath(src):
        return {"rel": rel, "action": action, "noop": True}
    if os.path.exists(dst):
        return {"error": f"目標已有同名檔：{name}"}
    set_decision(rel, "auto")                         # 清掉舊人工 decision（物理位置即決策）
    os.replace(src, dst)
    return {"rel": os.path.relpath(dst, _SRC), "from": rel, "action": action}


def _auto_decision(fn: str, ext: str) -> str:
    """自動判斷預設：乾淨副檔且檔名非逐字/字幕/校正 → extract；否則 exclude（§3 鐵律）。"""
    if ext in _SUPPORTED and not _DIRTY_NAME.search(fn):
        return "extract"
    return "exclude"


def _effective_decision(rel: str, fn: str, ext: str, decisions: dict) -> tuple[str, str | None]:
    """回 (生效決策, 人工覆寫值或None)。人工覆寫優先於自動判斷。"""
    override = decisions.get(rel)
    return (override or _auto_decision(fn, ext)), override

_CONF = {"高": 0.9, "中": 0.7, "低": 0.5}
_SRT_TS = re.compile(r"^\d+$|-->|^\s*$")
_BAD_FN = re.compile(r'[\\/:*?"<>|\s]+')
_SENT = re.compile(r"[。！？，、；：…]")          # 中文句讀＝整句修正，非術語
_ASCII_SENT = re.compile(r"[.!?;:,\"]")          # 英文句讀/逗號/引號＝填充語或整句（"All right." 漏網）
_PURE_NUM = re.compile(r"^[\d.,]+$")             # 純數字/量值（17.4、0.73）
_END_PUNCT = re.compile(r"[。！？!?．.、]$")
_LATIN_FILLER = {"all right", "well thank you", "i think so", "you know",
                 "ok", "okay", "right", "yeah", "thank you", "i mean"}
_SKIP_CAT = {"全句修正", "日語校正差異"}            # 對照表中的整句/翻譯校正列，非詞典詞

# 知識側只收 authored 乾淨來源（§3 鐵律）：逐字稿/字幕/ASR 校正含錯字，會污染知識側，一律排除。
_CLEAN_EXT = {".xlsx", ".docx", ".pptx", ".pdf"}        # 乾淨副檔；.srt/.txt（逐字/字幕）不收
_DIRTY_NAME = re.compile(r"字幕|逐字|校正|三語對齊|transcript|L3")  # 檔名含此＝逐字/ASR 校正來源，排除


def _is_term(s: str) -> bool:
    """術語守門：太長/含中英句讀/純數字/英文填充語＝非術語，丟（§3 知識側只收『詞』非『句』）。"""
    s = (s or "").strip()
    if not s or len(s) > 16:
        return False
    if _SENT.search(s) or _ASCII_SENT.search(s):          # 中/英標點＝整句或填充語（修 bug #1）
        return False
    if _PURE_NUM.match(s):                                  # 純數字/量值
        return False
    if s.lower() in _LATIN_FILLER:
        return False
    return True


def _clean_name(s: str, iso: str) -> str:
    """某語言『固定譯法』守門：按書寫系統限長度/詞數，擋整句（過長/句末標點/內嵌引號）。
    回乾淨譯法或空字串（修 bug #2，多語通用：CJK/泰文限字數、拼音文字限字數+詞數）。"""
    s = (s or "").strip().strip("「」『』\"'")
    if not s:
        return ""
    maxlen, maxwords = langtag.name_limits(iso)
    if len(s) > maxlen:
        return ""
    if _END_PUNCT.search(s) or _SENT.search(s):
        return ""
    if "「" in s or "」" in s or "。" in s:
        return ""
    if maxwords and len([w for w in s.split() if w]) > maxwords:   # 拼音文字超 N 詞＝句，非固定譯法
        return ""
    return s


# LLM 物件裡非語言的欄位（其餘頂層鍵若是註冊表語言代碼，即視為該語言固定譯法）。
_NON_LANG_KEYS = {"name", "term", "type", "domain", "meaning", "gloss", "desc",
                  "reason", "variants", "source", "target", "names", "src", "is_special", "occ"}

# 對照表序列化用的結構字（不得被當實體抽出，免污染心智圖）。
_GLOSSARY_STOP = {"正解", "誤聽變體", "誤聽", "變體", "配對", "對照表"}


def _parse_names(obj: dict) -> dict:
    """從 LLM 物件抽多語固定譯法 → {ISO-639-3: 乾淨值}（規格(一) 任意語系，加語言免改 code）。
    語言代碼『直接平鋪當頂層欄位』（en/ja/ko/...，保 JSON 扁平讓 _iter_objs 解得到）；
    亦相容 nested names；鍵經 langtag 正規化（未知鍵丟，不臆測），值按書寫系統清整。"""
    raw: dict = {}
    for k, v in obj.items():                               # 頂層平鋪語言欄位
        if k not in _NON_LANG_KEYS:
            raw[k] = v
    nm = obj.get("names")                                  # 相容 nested（多半被 _iter_objs 拆走，保險用）
    if isinstance(nm, dict):
        for k, v in nm.items():
            raw.setdefault(k, v)
    out: dict = {}
    for k, v in raw.items():
        if not isinstance(v, (str, int)):
            continue
        iso = langtag.canon_lang(k)
        if not iso:                                        # 不在註冊表的語言鍵 → 丟（避免污染）
            continue
        val = _clean_name(str(v), iso)
        if val:
            out[iso] = val
    return out


def _name_src(names: dict, text: str) -> dict:
    """每語言對應的來源（確定性，不問 LLM）：寫法出現在文本＝text(高信心)、否則＝gen(模型補,candidate)。
    （翻譯對照表的強制配對另標 glossary，在 _read_xlsx 設定。）"""
    return {iso: ("text" if v and str(v) in (text or "") else "gen") for iso, v in (names or {}).items()}


# ── 實體類型分類（§11）：心智圖/CAG 依命名實體種類分群（GraphRAG 標準實體分類）─────
# 8 類實體 + 未分類＝誠實兜底（不偽裝成類型，修 ③；取代舊「領域專名」catch-all）。
_UNCAT = "未分類"
_TAXO = ["人名", "組織機構", "地點地名", "產品品牌", "技術或概念術語",
         "外來詞或英文詞", "事件活動", "作品或標題"]
_TYPE_DESC = {
    "人名": "具體人物、姓名、綽號、職稱代稱",
    "組織機構": "公司、單位、部門、團體、政府機關、企業",
    "地點地名": "國家、城市、地區、場所、地標、位置",
    "產品品牌": "產品、服務、品牌、系統、平台的名稱",
    "技術或概念術語": "理論、方法、法則、原則、專業術語與具名概念",
    "外來詞或英文詞": "英文／外語詞、音譯詞、縮寫，需固定譯法的外來語",
    "事件活動": "會議、活動、專案、計畫、典禮、流程等具名事件",
    "作品或標題": "書名、章節、課程、文章、簡報、題目",
}
# few-shot：每類 1 個『通用常識』示範（刻意不舉本語料特例，避免 prompt 過擬合單一講者/領域）。
_TYPE_EXAMPLES = [
    ("愛因斯坦", "人名"), ("聯合國", "組織機構"),
    ("富士山", "地點地名"), ("iPhone", "產品品牌"),
    ("相對論", "技術或概念術語"), ("供需法則", "技術或概念術語"),
    ("KPI", "外來詞或英文詞"), ("奧林匹克運動會", "事件活動"),
    ("紅樓夢", "作品或標題"),
]
# 雜類別關鍵字 → 實體類型（LLM 多直接回標準名，此為容錯）。
_CAT_RULES = [
    ("人名", ("人名", "人物", "姓名", "綽號", "職稱")),
    ("組織機構", ("組織", "機構", "公司", "單位", "部門", "團體", "機關", "企業")),
    ("地點地名", ("地點", "地名", "國家", "城市", "地區", "場所", "地標", "位置")),
    ("產品品牌", ("產品", "品牌", "服務", "系統", "平台", "商品")),
    ("技術或概念術語", ("技術", "概念", "術語", "理論", "方法", "法則", "原則", "思維", "知識", "主題")),
    ("外來詞或英文詞", ("外來", "英文", "外語", "音譯", "縮寫", "abbr")),
    ("事件活動", ("事件", "活動", "會議", "專案", "計畫", "計劃", "典禮", "流程", "簽核")),
    ("作品或標題", ("作品", "標題", "書名", "章節", "課程", "文章", "簡報", "題目")),
]


def _taxo_block() -> str:
    """實體類型定義 + few-shot 範例（餵 Pass1/Pass2/tune 提示詞，幫 7B 用對視角分群）。"""
    defs = "；".join(f"{k}（{v}）" for k, v in _TYPE_DESC.items())
    ex = "；".join(f"{w}＝{d}" for w, d in _TYPE_EXAMPLES)
    return (f"實體類型限定（依命名實體種類分類，擇一）：{defs}。\n"
            f"判不準就標「{_UNCAT}」，不要硬塞。\n範例：{ex}。\n")


def _canon_category(cat: str) -> str:
    """雜類別 → 實體類型之一；無法判斷 → 未分類（誠實兜底，不偽裝成類型，修 ③）。"""
    c = (cat or "").strip()
    if not c or c == _UNCAT:
        return _UNCAT
    if c in _TAXO:
        return c
    for canon, kws in _CAT_RULES:
        if any(k in c for k in kws):
            return canon
    return _UNCAT


# ── 檔案解析 ─────────────────────────────────────────────
def _read_docx(path: str) -> str:
    import docx
    return "\n".join(p.text for p in docx.Document(path).paragraphs)


def _read_pptx(path: str) -> str:
    from pptx import Presentation
    out = []
    for slide in Presentation(path).slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
    return "\n".join(out)


def _read_pdf(path: str) -> str:
    from pypdf import PdfReader
    return "\n".join((pg.extract_text() or "") for pg in PdfReader(path).pages)


def _read_srt(path: str) -> str:
    with open(path, encoding="utf-8", errors="ignore") as f:
        return "\n".join(ln for ln in (l.strip() for l in f) if ln and not _SRT_TS.match(ln))


def _read_txt(path: str) -> str:
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()


# ── xlsx 對照表（錯字→正確字字典）智慧解析 ───────────────
def _col(headers: list, *keys: str):
    """以關鍵字找欄 index；避開計數欄（含『數』如 誤辨文字數/筆數）。回 None 表無此欄。"""
    for i, h in enumerate(headers):
        hs = str(h or "")
        if any(k in hs for k in keys) and "數" not in hs:
            return i
    return None


def _clean_variants(raw, correct: str, existing=None) -> list[str]:
    """變體清理（去空白/數字/非術語/去重/排除等於正解）——統一三處重複邏輯。"""
    cf = _flat(correct)
    out = list(existing or [])
    seen = {_flat(v) for v in out}
    for v in raw:
        v = (v or "").strip()
        if v and not v.isdigit() and _is_term(v) and _flat(v) != cf and _flat(v) not in seen:
            seen.add(_flat(v)); out.append(v)
    return out


def _read_xlsx(path: str) -> list[dict]:
    """解析對照表 xlsx → 統一 term 列。每分頁自動判型：
    - 校正型：有『正確/建議修正/正名』欄 → kind=correction，term + variants（錯字/辨識/誤辨/錯誤欄）。
    - 翻譯型：有語言欄、無正確欄 → kind=translation，names={ISO:寫法}（強制配對；任意 26 語，
      表頭經 langtag.canon_lang 認 zh/en/jp/中文/英語/Deutsch… 都可；錨語言取中文，無則第一個語言欄）。
    回 [{term, variants, names, kind, category, confidence, source, context, sheet}]（同 term 合併）。
    註：跨語『配對只保證正確/對照，非完美翻譯』——故一律標 src=glossary(鎖定)、不宣稱正規翻譯。"""
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    fn = os.path.splitext(os.path.basename(path))[0]
    merged: dict[str, dict] = {}
    for ws in wb.worksheets:
        rows = list(ws.iter_rows(values_only=True))
        hdr_i = next((i for i, r in enumerate(rows) if any(c is not None for c in r)), None)
        if hdr_i is None:
            continue
        header = list(rows[hdr_i])
        data = rows[hdr_i + 1:]                    # 跳過表頭列（否則「正確zh」「Zh」等表頭被當 term 抽出）
        sheet = (ws.title or "").strip()
        c_correct = _col(header, "正確", "建議修正", "正名")
        lang_cols = {i: iso for i, h in enumerate(header)          # 任意 26 語表頭 → ISO（langtag 通用）
                     if (iso := langtag.canon_lang(str(h or "").strip()))}
        if c_correct is not None:                     # ── 校正型 ──
            c_wrong = _col(header, "錯字", "辨識", "誤辨", "錯誤")
            c_type, c_conf = _col(header, "類型", "分類"), _col(header, "信心")
            c_ctx = _col(header, "上下文", "範例", "彙總")
            for r in data:
                correct = str(r[c_correct]).strip() if c_correct < len(r) and r[c_correct] else ""
                if not _is_term(correct):
                    continue
                cat = str(r[c_type]).strip()[:16] if (c_type is not None and c_type < len(r) and r[c_type]) else ""
                if cat in _SKIP_CAT:
                    continue
                e = merged.setdefault(_flat(correct), {"term": correct, "variants": [], "names": {}, "src": {},
                    "kind": "correction", "category": "", "confidence": 0.0,
                    "source": f"glossary:{fn}", "context": "", "sheet": sheet})
                if c_wrong is not None and c_wrong < len(r) and r[c_wrong]:
                    e["variants"] = _clean_variants(re.split(r"[、,，;；/\s]+", str(r[c_wrong])), correct, e["variants"])
                if cat and not e["category"]:
                    e["category"] = cat
                conf = _CONF.get(str(r[c_conf]).strip(), 0.8) if (c_conf is not None and c_conf < len(r) and r[c_conf]) else 0.8
                e["confidence"] = max(e["confidence"], conf)
                if c_ctx is not None and c_ctx < len(r) and r[c_ctx] and not e["context"]:
                    e["context"] = str(r[c_ctx]).strip()[:80]
        elif lang_cols:                               # ── 翻譯型（強制配對；任意 26 語）──
            c_anchor = next((i for i, iso in lang_cols.items() if iso == "cmn"), min(lang_cols))
            for r in data:
                term = str(r[c_anchor]).strip() if c_anchor < len(r) and r[c_anchor] else ""
                if not _is_term(term):
                    continue
                e = merged.setdefault(_flat(term), {"term": term, "variants": [], "names": {}, "src": {},
                    "kind": "translation", "category": "", "confidence": 0.9,
                    "source": f"glossary:{fn}", "context": "", "sheet": sheet})
                for i, iso in lang_cols.items():
                    if i != c_anchor and i < len(r) and r[i]:
                        val = _clean_name(str(r[i]).strip(), iso)
                        if val and not e["names"].get(iso):
                            e["names"][iso] = val
                            e["src"][iso] = "glossary"          # 強制配對＝鎖定來源（保證正確/對照，非完美翻譯）
    wb.close()
    return list(merged.values())


async def _glossary_hint(rows: list[dict], title: str) -> str:
    """一句廣泛描述對照表的資料性質/領域（per-table 一次 LLM，幾萬列也只一次）。
    當豐化前綴讓冷僻詞也有語意框架可掛 → type 判得準、gloss 寫得出、少落未分類。"""
    sample = "、".join(r["term"] for r in rows[:40] if r.get("term"))
    try:
        raw = await clients.judge_chat(
            f"對照表標題「{title}」，部分詞：{sample}。\n"
            "用一句話廣泛描述這是什麼領域、什麼類型的資料（供理解每個詞的語意）。只回那一句、勿列點。",
            max_tokens=150)
        return (raw or "").strip().replace("\n", " ")[:150]
    except Exception:
        return ""


def _serialize_glossary(rows: list[dict], fn: str, hint: str = "") -> str:
    """對照表 → 帶標題文本（餵 Gemma 補語意）。標題＝檔名＋分頁名 + 資料性質前綴(hint)，
    讓 Gemma 抓大致語意/領域；正解/變體/配對標鎖定、誤聽變體明示不立節點。"""
    by_sheet: dict[str, list[dict]] = {}
    for r in rows:
        by_sheet.setdefault(r.get("sheet") or fn, []).append(r)
    blocks = []
    for sheet, rs in by_sheet.items():
        head = (f"【對照表：{fn}／{sheet}】" + (f"資料性質：{hint}\n" if hint else "") +
                "此表正解、誤聽變體、跨語配對均已人工鎖定，請勿改動；"
                "請依上述資料性質，為每個『正解』補 gloss(一句語意，即使是『屬於某領域的某類事物』這種廣泛定位也要寫)、"
                "type(八大類)、與其他正解的關係；『誤聽變體』是 ASR 錯字、不要當實體抽出。\n")
        lines = []
        for r in rs:
            seg = "正解：" + r["term"]
            if r.get("names"):
                seg += "｜配對：" + "、".join(f"{k}={v}" for k, v in r["names"].items())
            if r.get("variants"):
                seg += "｜誤聽變體：" + "、".join(r["variants"])
            lines.append(seg)
        blocks.append(head + "\n".join(lines))
    return "\n\n".join(blocks)


_TEXT_EXT = {".docx": _read_docx, ".pptx": _read_pptx, ".pdf": _read_pdf,
             ".srt": _read_srt, ".txt": _read_txt}


# ── 切 chunk + LLM 抽詞+語意（精度層，取代統計新詞）─────────
# 整份餵門檻（字元）：≤此值的乾淨檔一次餵給判官（全域脈絡）——讓 Gemma 先懂整個檔在講什麼，
# 才判得出「標題/作品」這類需全域脈絡的型別、並抓跨段語意；超過才退回貪婪打包（避 lost-in-the-middle）。
# 貼齊 BUILD_MAX_MODEL_LEN（CJK 約 1.3–1.7 token/字，留指令+輸出餘裕）。
# 8192 context 固定：整份餵門檻壓在 3000 字（CJK 約 1.3–1.7 token/字，留指令+2400 輸出餘裕，不破 8192）。
_WHOLE_DOC_MAX = 3000


def _doc_title(text: str) -> str:
    """文件標題＝首個非空行（docx 首段／pptx 首頁標題）。長檔分塊時夾帶給判官當全域錨。"""
    for ln in (text or "").splitlines():
        ln = ln.strip()
        if ln:
            return ln[:60]
    return ""


def _chunk_text(text: str, size: int = 700, title: str | None = None) -> list[str]:
    """乾淨文件切 chunk（句界貪婪打包，保留英文/混語）。長檔才用；短檔走整份餵（見 _extract_doc）。
    title：分塊時每塊夾帶【文件標題】——讓判官即使只看到某塊也知道整檔在講什麼（才判得出標題/作品、不誤判內文）。"""
    parts, buf, n = [], [], 0
    for seg in re.split(r"(?<=[。！？!?\n])", text or ""):
        if not seg.strip():
            continue
        if buf and n + len(seg) > size:
            parts.append("".join(buf)); buf, n = [], 0
        buf.append(seg); n += len(seg)
    if buf:
        parts.append("".join(buf))
    if title:
        head = f"【文件標題】{title}\n"
        parts = [head + p for p in parts]
    return parts


# 兩階段抽取（coarse-to-fine）：Pass1 語意心智圖（廣度）→ Pass2 吃 Pass1 全局理解抽特殊詞 CAG（精度）。
# 分類軸＝命名實體類型（_TAXO：8 類實體，GraphRAG 標準分類），未分類誠實兜底（修 ④）。


def _iter_objs(raw: str):
    """物件解析（平衡括號，支援巢狀 names/src）：掃出頂層 {...}，json.loads；
    若含 entities/relationships 陣列就攤平 yield 每個元素，否則 yield 該物件本身。
    （26B 後 Pass1 改巢狀多語結構 names:{ISO:寫法}+src:{ISO:text/gen}，需平衡括號才解得到。）"""
    depth, start = 0, -1
    for i, ch in enumerate(raw):
        if ch == "{":
            if depth == 0:
                start = i
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0 and start >= 0:
                try:
                    obj = json.loads(raw[start:i + 1])
                except Exception:
                    start = -1
                    continue
                start = -1
                if isinstance(obj, dict) and (obj.get("terms") or obj.get("entities") or obj.get("relationships")):
                    for e in (obj.get("terms") or obj.get("entities") or []):
                        if isinstance(e, dict):
                            yield e
                    for r in (obj.get("relationships") or []):
                        if isinstance(r, dict):
                            yield r
                elif isinstance(obj, dict):
                    yield obj


# ── Pass 1：語意心智圖（type/gloss/多語/關係；廣度、召回）──────
_LANG_LINE = "、".join(f"{iso}={m['name']}" for iso, m in langtag.LANGS.items())


def _pass1_prompt(chunk: str, already: list[str] | None = None) -> str:
    """Pass1 語意心智圖（benchmark ③ 勝出版，2026-06-24）：語意廣度抽取 + 八大分類 + 跨語配對。
    names 用 ISO-639-3 巢狀 map、src 標 text/gen（gen＝模型生成譯法＝candidate，治理鐵律）。"""
    p = (
        "你是譯前語意分析師，為文本建立『語意心智圖』。通讀文本，抽出所有承載具體意義的實體"
        "（人事物、概念、專名），理解其語意並標跨語對應，供後續正確翻譯。\n"
        f"-八大類型- 擇一分類：{'、'.join(_TAXO)}。判不準標「未分類」。\n"
        f"-語言代碼(ISO-639-3)- {_LANG_LINE}。\n"
        "-每個實體輸出欄位-\n"
        "  term：實體在文本的原文寫法。\n"
        "  type：八大類型之一。\n"
        "  gloss：在此語境的真正意思，一句話（幫助跨語翻譯，非字典定義）。\n"
        "  names：跨語對應寫法 map，key 用上列 ISO 代碼、value 為該語言寫法；"
        "只填文本中實際出現、或你有十足把握的對應，沒把握的語言一律省略，嚴禁臆造。\n"
        "  src：對 names 每個條目標來源，'text'=文本中出現、'gen'=你依把握補的對應。\n"
        "並抽實體間的關係：source、target（皆須為上面抽出的實體 term）、desc（關係，如 屬於/部分/因果/提出）。\n"
        "-原則- 抽承載具體意義的實詞（廣度優先，語意完整）；略過虛詞與口語填充。不必判斷是否值得進字典。\n"
        '-輸出- 只輸出 JSON：{"terms":[{"term":"","type":"","gloss":"","names":{},"src":{}}],'
        '"relationships":[{"source":"","target":"","desc":""}]}\n')
    if already:
        p += ("\n本文本已抽出下列詞，請【只補出遺漏】的，不要重複：\n" + "、".join(already[:80]) + "\n")
    return p + "\n文本：\n" + chunk


async def _pass1_semantic(chunk: str, gleanings: int = 1) -> dict:
    """單 chunk 語意抽取（首趟 + gleaning 補抽）。回 {entities:{flat:e}, relationships:[]}。"""
    ents: dict[str, dict] = {}
    rels: list[dict] = []
    for i in range(1 + max(0, gleanings)):
        found = [e["term"] for e in ents.values()] if i else None
        try:
            raw = await clients.judge_chat(_pass1_prompt(chunk, found), max_tokens=2400)
        except Exception:
            break
        new = 0
        for obj in _iter_objs(raw):
            if obj.get("term") or obj.get("name"):
                n = str(obj.get("term") or obj.get("name")).strip()
                if n and not n.isdigit() and _is_term(n) and n in chunk and _flat(n) not in ents:
                    nm = _parse_names(obj)
                    ents[_flat(n)] = {"term": n,
                                      "type": str(obj.get("type") or obj.get("domain", "")).strip()[:12],
                                      "gloss": str(obj.get("gloss") or obj.get("meaning") or obj.get("desc", "")).strip()[:60],
                                      "names": nm, "src": _name_src(nm, chunk)}
                    new += 1
            elif obj.get("source") and obj.get("target"):
                s, t = str(obj["source"]).strip(), str(obj["target"]).strip()
                if s and t and _flat(s) != _flat(t):
                    rels.append({"source": s, "target": t, "desc": str(obj.get("desc", "")).strip()[:50]})
        if i and new == 0:
            break
    return {"entities": ents, "relationships": rels}


# ── Pass 2：特殊詞 CAG（吃 Pass1 全局理解；精度、不能錯）──────
# 術語表收錄判準（termhood + 翻譯關鍵性 + ASR 風險）；可操作、取代主觀「coined」。
_REASON = {"A": "命名實體", "B": "概念術語", "C": "易誤判"}


def _pass2_prompt(chunk: str, known: list[str], title: str = "") -> str:
    """輸入＝本文主題/領域(title) + Pass1 全局理解(known) + chunk；以 A/B/C 判準挑『值得進術語表』的詞。
    領域感知：同一個詞在其所屬領域是否該收，依本文主題判斷（地質文件的『海蝕』值得收，泛文不必）。"""
    dom = (f"本文主題／領域：{title}\n（請依此領域判斷專詞值不值得收——領域核心專詞要收、與領域無關的泛詞不收。）\n") if title else ""
    ctx = ("全篇關鍵詞脈絡：" + "、".join(known[:60]) + "\n") if known else ""
    # 文本在前、輸出指令在後（避免 7B 把長指令後的文本直接續寫、不產 JSON）。
    return (
        "你是為『跨語言翻譯』建立術語表的嚴格編纂員，寧缺勿濫——"
        "只收『翻譯時必須前後一致、譯錯會出事』的詞，供之後各語言固定譯法對齊。\n" + dom + ctx +
        "\n文本：\n" + chunk +
        "\n\n任務：從上面文本挑出需要進『翻譯術語表』的詞，符合下列其一才收：\n"
        "  A 專有名稱：人／地／組織／產品／品牌／作品等具名實體。\n"
        "  B 具名概念：有專屬名稱、不可逐字直譯的理論／方法／法則／口號（具名概念，非一般描述）。\n"
        "  C 同音或形近易誤辨、譯者易選錯字的專詞。\n"
        "鐵則：受過教育的譯者不需查資料就能譯對的一般詞（泛指的日常用語）一律不收。\n"
        + _taxo_block() +
        f"語言代碼(ISO-639-3)：{_LANG_LINE}。\n"
        "請輸出 JSON 陣列（沒有符合就回 []），term 照原文寫法、type 為上列實體類型、"
        "names＝跨語對應 map(key 用上列 ISO 代碼、value 為該語言寫法，有把握才給、整句不算)、"
        "src＝對 names 每條目標 text(文本出現)/gen(你補的)：\n"
        '[{"term":"","reason":"A","type":"","names":{},"src":{},"variants":[]}]')


async def _pass2_special(chunk: str, known: list[str], title: str = "") -> dict:
    """單 chunk 術語表收錄（吃全局 known + 領域 title，A/B/C 判準）。回 {flat: {term,reason,names,variants}}。"""
    try:
        raw = await clients.judge_chat(_pass2_prompt(chunk, known, title), max_tokens=4096)
    except Exception:
        return {}
    out: dict[str, dict] = {}
    for obj in _iter_objs(raw):
        t = str(obj.get("term", "")).strip()
        if not t or t.isdigit() or not _is_term(t) or t not in chunk:
            continue
        reason = str(obj.get("reason", "")).strip().upper()[:1]
        if reason not in _REASON:                  # 未給有效判準 → 不收（守 termhood）
            continue
        vs = obj.get("variants") or []
        variants = [str(v).strip() for v in vs
                    if isinstance(v, (str, int)) and _is_term(str(v).strip())
                    and not str(v).strip().isdigit() and _flat(str(v).strip()) != _flat(t)]
        nm = _parse_names(obj)
        out[_flat(t)] = {"term": t, "reason": reason, "variants": variants,
                         "type": str(obj.get("type") or obj.get("domain", "")).strip()[:12],
                         "names": nm, "src": _name_src(nm, chunk)}
    return out


async def _extract_doc(text: str, cap: int = 40, conc: int = 8, gleanings: int = 1, title: str = "") -> dict:
    """兩階段：Pass1 逐 chunk 語意（並行）→ 聚合文件級理解 → Pass2 逐 chunk 抽特殊詞（吃全局+領域 title）。
    title＝本文主題/領域（glossary 傳 hint、doc 自動取首行），餵 Pass2 做領域感知收錄。
    回 {entities:{flat:e}, relationships:[]}；entity 含語意面(type/gloss/names)＋正確度面(reason/variants/is_special)。"""
    dtitle = title or _doc_title(text)
    # 整份餵：短檔不切，讓判官讀懂整個檔在講什麼（才正確判標題/作品、抓跨段語意）；
    # 長檔退回貪婪打包，但每塊夾帶【文件標題】保留全域錨（避 lost-in-the-middle 又不失主題）。
    chunks = ([text] if 0 < len(text) <= _WHOLE_DOC_MAX
              else _chunk_text(text, size=1100, title=dtitle))
    chunks = chunks[:cap]
    if not chunks:
        return {"entities": {}, "relationships": []}
    sem = asyncio.Semaphore(conc)

    async def p1(ch):
        async with sem:
            return await _pass1_semantic(ch, gleanings)

    # Pass 1：語意（並行）→ 聚合
    res1 = await asyncio.gather(*[p1(c) for c in chunks])
    ents: dict[str, dict] = {}
    rels: list[dict] = []
    for r in res1:
        for k, e in r["entities"].items():
            cur = ents.setdefault(k, {"term": e["term"], "type": e["type"], "gloss": e["gloss"],
                                      "names": dict(e.get("names") or {}), "src": dict(e.get("src") or {}),
                                      "occ": 0, "reason": "", "variants": [], "is_special": False})
            cur["occ"] += 1
            for f in ("gloss", "type"):
                if e[f] and not cur[f]:
                    cur[f] = e[f]
            for lang, val in (e.get("names") or {}).items():        # 多語固定譯法+來源：跨 chunk 只補空
                if val and not cur["names"].get(lang):
                    cur["names"][lang] = val
                    cur["src"][lang] = (e.get("src") or {}).get(lang, "gen")
        rels += r["relationships"]

    # Pass 2：術語表收錄（吃 Pass1 全局理解，A/B/C 判準）
    known = [e["term"] for e in sorted(ents.values(), key=lambda x: -x["occ"])]

    async def p2(ch):
        async with sem:
            return await _pass2_special(ch, known, dtitle)

    res2 = await asyncio.gather(*[p2(c) for c in chunks])
    for sp in res2:
        for k, s in sp.items():
            cur = ents.get(k)
            if not cur:                       # Pass2 撈到 Pass1 漏的收錄詞 → 補成實體（未分類兜底，修 ①）
                cur = ents[k] = {"term": s["term"], "type": s.get("type") or _UNCAT, "gloss": "",
                                 "names": {}, "src": {}, "occ": 1, "reason": "",
                                 "variants": [], "is_special": False}
            cur["is_special"] = True
            cur["reason"] = cur["reason"] or s["reason"]
            if s.get("type"):                 # Pass2 吃全局脈絡定類型 → 覆寫 Pass1 局部判斷（綜合判斷，修 ①）
                cur["type"] = s["type"]
            cur.setdefault("src", {})
            for lang, val in (s.get("names") or {}).items():        # Pass2 多語固定譯法+來源只補空
                if val and not cur["names"].get(lang):
                    cur["src"][lang] = (s.get("src") or {}).get(lang, "gen")
                    cur["names"][lang] = val
            for v in s["variants"]:
                if v not in cur["variants"]:
                    cur["variants"].append(v)
    return {"entities": ents, "relationships": rels}


# ── 單檔抽取（產可序列化記錄，供 manifest 快取）──────────────
async def _extract_file(path: str, fn: str, ext: str,
                        use_llm: bool, llm_cap: int, gleanings: int, force: bool = False) -> dict:
    """一個來源檔 → 記錄 {kind, ...}。kind=glossary/doc/skip。doc/glossary 內容可 JSON 快取。
    force=True（使用者手動指定抽取）→ 略過逐字/字幕/校正與非乾淨副檔的自動排除門。"""
    if _DIRTY_NAME.search(fn) and not force:            # 逐字/字幕/校正/三語對齊：含 ASR 錯字，排除（§3 鐵律）
        return {"kind": "skip", "note": fn + "(逐字/字幕/校正，排除)"}
    if ext == ".xlsx":
        rows = _read_xlsx(path)
        if not rows:
            return {"kind": "skip", "note": fn + "(非對照表)"}
        ents, rels = {}, []
        if use_llm:                                     # 也給 Gemma：序列化帶標題+資料性質前綴 → 補 gloss/type/關係
            base = os.path.splitext(fn)[0]
            hint = await _glossary_hint(rows, base)
            block = _serialize_glossary(rows, base, hint)
            doc = await _extract_doc(block, cap=llm_cap, gleanings=gleanings, title=hint or base)
            vk = {_flat(v) for r in rows for v in r.get("variants", [])}   # 誤聽變體不立節點
            ents = {k: e for k, e in doc["entities"].items()
                    if k not in vk and e["term"] not in _GLOSSARY_STOP}    # 並擋序列化結構字
            rels = doc["relationships"]
        return {"kind": "glossary", "rows": rows, "entities": ents, "relationships": rels,
                "note": f"{fn}→{len(rows)}詞(glossary)+{len(ents)}語意"}
    if ext in _TEXT_EXT:
        if ext not in _CLEAN_EXT and not force:         # .srt/.txt＝逐字稿/字幕來源，知識側不收（force 可強收）
            return {"kind": "skip", "note": fn + f"({ext} 逐字/字幕來源，排除)"}
        if not use_llm:
            return {"kind": "skip", "note": fn + "(use_llm=False)"}
        doc = await _extract_doc(_TEXT_EXT[ext](path), cap=llm_cap, gleanings=gleanings)
        sp = sum(1 for e in doc["entities"].values() if e["is_special"])
        return {"kind": "doc", "entities": doc["entities"], "relationships": doc["relationships"],
                "note": f"{fn}→{len(doc['entities'])}實體/{sp}特殊/{len(doc['relationships'])}關係"}
    return {"kind": "skip", "note": f"{fn}({ext or '無副檔'})"}


def _aggregate(rec: dict, src: str, terms: dict, rels: list) -> None:
    """把單檔抽取記錄併入 terms/rels（增量：跨檔累加 occ、union variants/names/關係）。"""
    if rec["kind"] == "glossary":
        for e in rec["rows"]:
            k = _flat(e["term"])
            cur = terms.get(k)
            if not cur or e["confidence"] >= cur.get("confidence", 0):
                terms[k] = {**(cur or {}), **e,
                            "variants": _clean_variants(e["variants"], e["term"], (cur or {}).get("variants", [])),
                            "names": {**(cur or {}).get("names", {}), **e.get("names", {})},      # 強制配對保留
                            "src": {**(cur or {}).get("src", {}), **e.get("src", {})}}            # glossary 來源標記
        # Gemma 補的語意回填 gold（只補空 gloss/類別/語言，不碰鎖定的正解/變體/配對）
        for k, ent in rec.get("entities", {}).items():
            cur = terms.get(k)
            if not cur:                                   # glossary 文本只豐化已知 gold，不新增節點
                continue
            if ent.get("gloss") and not cur.get("gloss"):
                cur["gloss"] = ent["gloss"]
            if ent.get("type") and not cur.get("category"):
                cur["category"] = ent["type"]
            for lang, val in (ent.get("names") or {}).items():   # 補空語言(gen)，不覆蓋鎖定配對
                if val and not cur.setdefault("names", {}).get(lang):
                    cur["names"][lang] = val
                    cur.setdefault("src", {})[lang] = (ent.get("src") or {}).get(lang, "gen")
        rels.extend(rec.get("relationships", []))
    elif rec["kind"] == "doc":
        for k, e in rec["entities"].items():
            cur = terms.get(k)
            # glossary gold：term/變體/類別/信心/來源不動，但開放心智圖補『空的』語意欄（gloss/names/關係）——
            # 第一原則（語意接地）貫穿到對照表詞，且只補空不覆寫，不污染人工核可資料。
            gold = bool(cur and cur.get("source", "").startswith("glossary"))
            if not cur:
                terms[k] = cur = {"term": e["term"], "variants": [], "category": e.get("type") or _UNCAT,
                                  "confidence": 0.0, "source": f"llm:{src}", "context": "",
                                  "gloss": e.get("gloss", ""), "occ": 0, "files": set(),
                                  "names": dict(e.get("names") or {}), "src": dict(e.get("src") or {}),
                                  "reason": "", "is_special": False}
            cur["occ"] = cur.get("occ", 0) + e.get("occ", 1)
            cur.setdefault("files", set()).add(src)
            if e.get("gloss") and not cur.get("gloss"):
                cur["gloss"] = e["gloss"]
            for lang, val in (e.get("names") or {}).items():     # 多語固定譯法+來源只補空（gold 也補：第一原則）
                if val and not cur.setdefault("names", {}).get(lang):
                    cur["names"][lang] = val
                    cur.setdefault("src", {})[lang] = (e.get("src") or {}).get(lang, "gen")
            cur["reason"] = cur.get("reason") or e.get("reason")
            cur["is_special"] = cur.get("is_special") or e.get("is_special")
            if not gold:                                          # 非 gold 才併變體（保護人工核可的 ASR 變體不被 LLM 污染）
                for v in e.get("variants", []):
                    if v not in cur["variants"]:
                        cur["variants"].append(v)
        rels.extend(rec["relationships"])


# ── Tune 階段（§3）：心智圖出來後，對著全局重判落不到實體類型的詞 ──────────
async def _tune_categories(terms: dict[str, dict], conc: int = 8) -> dict:
    """心智圖先成形 → 把『未分類/不落 _TAXO』的詞對著**整張圖**重判實體類型（綜合判斷）。
    只 tune misfit（便宜先跑、gold/已落類型不動），用全局已分群詞當 few-shot 參照。
    就地更新 terms[k]['category']。回 {tuned, recovered}。"""
    misfit = [k for k, e in terms.items() if (e.get("category") or "") not in _TAXO]
    if not misfit:
        return {"tuned": 0, "recovered": 0}
    classified: dict[str, list[str]] = {}            # 全局已落場景的詞＝資料自身 few-shot 參照
    for e in terms.values():
        c = e.get("category")
        if c in _TAXO:
            classified.setdefault(c, []).append(e["term"])
    ctx = "；".join(f"{c}：{('、'.join(ts[:12]))}" for c, ts in classified.items())
    sem = asyncio.Semaphore(conc)
    recovered = 0

    async def tune_batch(batch: list[str]) -> int:
        nonlocal recovered
        async with sem:
            items = "、".join(
                f"{terms[k]['term']}（{terms[k].get('gloss') or '—'}）" for k in batch)
            prompt = (
                "你在整理一張命名實體知識圖"
                "（人名/組織機構/地點地名/產品品牌/技術或概念術語/外來詞或英文詞/事件活動/作品或標題）。\n"
                "全局已分群參照：\n" + (ctx or "（尚無）") + "\n"
                + _taxo_block() +
                "請依命名實體種類把下列每個詞歸到一個實體類型（真的判不出才標「" + _UNCAT + "」）。\n"
                "詞（附語意）：" + items + "\n"
                '只輸出 JSON 陣列：[{"term":"","type":""}]')
            try:
                raw = await clients.judge_chat(prompt, max_tokens=900)
            except Exception:
                return 0
            for obj in _iter_objs(raw):
                t = str(obj.get("term", "")).strip()
                d = _canon_category(str(obj.get("type") or obj.get("domain", "")).strip())
                k = _flat(t)
                if k in terms and d != _UNCAT and (terms[k].get("category") or "") not in _TAXO:
                    terms[k]["category"] = d
                    recovered += 1
            return 0

    batches = [misfit[i:i + 20] for i in range(0, len(misfit), 20)]
    await asyncio.gather(*[tune_batch(b) for b in batches])
    return {"tuned": len(misfit), "recovered": recovered}


# ── 萃取彙總（增量 + 證據分層）──────────────────────────────
async def _harvest_sources(use_llm: bool = True, llm_cap: int = 40, gleanings: int = 1,
                           incremental: bool = True) -> dict:
    """遞迴走 sources/ → 內容 hash 增量抽取（只重抽新/改動檔，其餘讀 manifest 快取）→ 合併。
    回 {terms,relationships,used,skipped,reused}。"""
    manifest = _load_manifest() if incremental else {}
    decisions = _load_decisions()                  # 每檔人工抽取/排除覆寫
    new_manifest: dict = {}
    terms: dict[str, dict] = {}
    rels: list[dict] = []
    skipped, used, reused = [], [], 0
    for root, _dirs, files in os.walk(_SRC):       # 遞迴：支援 sources/<領域>/ 子目錄
        _dirs[:] = [d for d in _dirs if not d.startswith(("_", ".")) and d != "package"]  # 略過 _excluded/封存、package(檢視產物)
        for fn in sorted(files):
            if fn == ".gitkeep":
                continue
            path = os.path.join(root, fn)
            rel = os.path.relpath(path, _SRC)
            ext = os.path.splitext(fn)[1].lower()
            eff, override = _effective_decision(rel, fn, ext, decisions)
            if eff == "exclude":                       # 自動判斷或人工指定排除 → 不抽
                tag = "人工排除" if override else "自動排除"
                skipped.append(f"{fn}（{tag}）")
                continue
            force = override == "extract"              # 人工指定抽取 → 強收（略過自動排除門）
            try:
                h = _file_hash(path)
                cached = manifest.get(rel)
                hit = bool(incremental and cached and cached.get("hash") == h)
                if hit:
                    rec = cached                       # 內容未變 → 讀快取，免重抽
                    reused += 1
                else:
                    rec = await _extract_file(path, fn, ext, use_llm, llm_cap, gleanings, force=force)
                    rec["hash"] = h
            except Exception as ex:
                skipped.append(f"{fn}(錯誤:{str(ex)[:40]})")
                continue
            if rec["kind"] == "skip":
                skipped.append(rec["note"])            # skip 不快取（便宜、每次重判，避免 use_llm 切換誤用）
                continue
            new_manifest[rel] = rec                     # 只快取 glossary/doc（昂貴的抽取結果）
            used.append(rec["note"] + (" [快取]" if (cached and cached.get("hash") == h) else ""))
            _aggregate(rec, rel, terms, rels)
    _save_manifest(new_manifest)        # 永遠存：incremental 只控「讀不讀舊快取」，全抽後仍須存以初始化增量

    # 證據分層（兩階段）：glossary=T1(gold)；LLM 實體中 Pass2 標特殊(is_special)=T2(CAG 候選)、純語意=T3(僅心智圖)
    for e in terms.values():
        if e.get("source", "").startswith("glossary"):
            e["tier"] = "T1" if e.get("variants") else "T1g"
            e["is_special"] = True                            # 對照表本就是收錄詞 → 進 CAG
            e.setdefault("reason", "") or e.update(reason="A")  # 對照表多為命名實體/專名
            e["category"] = _canon_category(e.get("category"))  # 對照表 類型欄 → 實體類型正規化
        else:
            e["tier"] = "T2" if e.get("is_special") else "T3"
            e["confidence"] = 0.7 if e["tier"] == "T2" else 0.5
            e["category"] = _canon_category(e.get("category"))  # 實體類型正規化（8 類 + 未分類）
        if isinstance(e.get("files"), set):
            e["files"] = sorted(e["files"])

    # 關係去重 + 只留兩端皆為已收實體者（連得到真節點）
    rel_map: dict[tuple, dict] = {}
    for r in rels:
        ks, kt = _flat(r["source"]), _flat(r["target"])
        if ks in terms and kt in terms and ks != kt:
            rel_map.setdefault((ks, kt), {"source": terms[ks]["term"],
                                          "target": terms[kt]["term"], "desc": r["desc"]})
    return {"terms": terms, "relationships": list(rel_map.values()),
            "used": used, "skipped": skipped, "reused": reused}


# ── Obsidian vault + CAG 寫出 ────────────────────────────
def _safe(name: str) -> str:
    return _BAD_FN.sub("_", name).strip("_")[:80] or "term"


def _write_vault(terms: dict[str, dict], relationships: list[dict] | None = None) -> dict:
    os.makedirs(os.path.join(_VAULT, "terms"), exist_ok=True)
    os.makedirs(os.path.join(_VAULT, "domains"), exist_ok=True)
    # 同音群（IPA key → terms）做 [[同音]] 連結（2026-06-24：拼音 key 改 IPA，全專案統一音韻表徵）
    by_key: dict[str, list[str]] = {}
    for e in terms.values():
        e["ipa"] = to_ipa(e["term"])
        if e["ipa"]:
            by_key.setdefault(e["ipa"], []).append(e["term"])
    # 關係邊：source term → [(target, desc)]（GraphRAG 語意邊，餵心智圖）
    rel_by: dict[str, list[tuple]] = {}
    for r in (relationships or []):
        rel_by.setdefault(r["source"], []).append((r["target"], r.get("desc", "")))
    domains: dict[str, list[str]] = {}
    for e in terms.values():
        # 鎖定欄位保護（§11）：人工核可(locked:true)的 note，語意欄位沿用人工版；
        # build 仍累加 occ/關係/tier，但不覆寫人工改過的 category/gloss/variants/reason/names。
        notep = os.path.join(_VAULT, "terms", _safe(e["term"]) + ".md")
        locked = False
        if os.path.exists(notep):
            old, _r = _parse_term_note(notep)
            if old.get("locked"):
                locked = True
                for f in ("category", "gloss", "reason"):
                    if old.get(f) not in (None, ""):
                        e[f] = old[f]
                # 變體無條件沿用人工版（含清空 []）——人工清掉的髒變體 build 不得從對照表 union 還原
                e["variants"] = old.get("variants") or []
                if old.get("names"):
                    e["names"] = old["names"]
                if old.get("names_src"):
                    e["src"] = old["names_src"]
                if old.get("term_lang"):
                    e["term_lang"] = old["term_lang"]
        e["_locked"] = locked
        cat = e.get("category") or "未分類"
        domains.setdefault(cat, []).append(e["term"])
        homo = [t for t in by_key.get(e["ipa"], []) if t != e["term"]]
        names = _names_map(e)                          # 多語對照（ISO-639-3 key，可擴充）
        src = e.get("src") or {}                        # 每語言來源：text(文本)/gen(模型)/glossary(強制配對)
        term_lang = e.get("term_lang") or langtag.lang_tag(e["term"])
        names_str = ", ".join(f"{k}: {v}" for k, v in names.items())
        src_str = ", ".join(f"{k}: {src[k]}" for k in names if k in src)
        fm = [
            "---", f"term: {e['term']}", f"term_lang: {term_lang}",
            f"ipa: {e['ipa']}",
            f"category: {cat}", f"source: {e['source']}",
            f"tier: {e.get('tier', '')}", f"confidence: {e['confidence']}",
            f"reason: {e.get('reason', '')}",
            f"is_special: {str(bool(e.get('is_special'))).lower()}",
            f"locked: {str(bool(e.get('_locked'))).lower()}",   # 人工核可後設 true → build 不覆寫語意欄位
            f"names: {{{names_str}}}",                  # 任意語言可擴充（DaMuEL 模式）
            f"names_src: {{{src_str}}}",                # 對應每語言來源 text/gen/glossary（治理：gen=candidate）
            f"variants: [{', '.join(e['variants'])}]", "---", "",
            f"# {e['term']}", "",
        ]
        if e.get("gloss"):
            fm += [f"*{e['gloss']}*", ""]
        if names:
            fm += ["多語： " + " / ".join(f"{k}: {v}" for k, v in names.items()), ""]
        if e.get("context"):
            fm += [f"> {e['context']}", ""]
        if e["variants"]:
            fm.append("ASR 錯法： " + "、".join(f"[[{v}]]" for v in e["variants"]))
        if homo:
            fm.append("同音： " + "、".join(f"[[{t}]]" for t in homo))
        for tgt, desc in rel_by.get(e["term"], []):       # GraphRAG 關係邊
            fm.append(f"關係： [[{tgt}]]" + (f"（{desc}）" if desc else ""))
        fm.append(f"領域： [[domains/{_safe(cat)}]]")
        with open(os.path.join(_VAULT, "terms", _safe(e["term"]) + ".md"), "w", encoding="utf-8") as f:
            f.write("\n".join(fm) + "\n")
    for cat, ts in domains.items():
        body = ["---", f"domain: {cat}", f"term_count: {len(ts)}", "---", "",
                f"# {cat}", "", "： ".join(["術語"]) + "： " + "、".join(f"[[{t}]]" for t in sorted(set(ts)))]
        with open(os.path.join(_VAULT, "domains", _safe(cat) + ".md"), "w", encoding="utf-8") as f:
            f.write("\n".join(body) + "\n")
    return {"notes": len(terms), "domains": len(domains)}


def _write_cag(terms: dict[str, dict], domain: str) -> dict:
    """CAG＝心智圖的精度子集：is_special 的詞（對照表 gold + Pass2 A/B/C 收錄），純語意 T3 排除。
    B 概念術語（需固定譯法、最不能直譯錯）優先列出。"""
    os.makedirs(_CAG, exist_ok=True)
    groups: dict[str, list[str]] = {}
    concept: list[str] = []                # reason B：非組合性概念術語（最該記住、最不能錯）
    for e in terms.values():
        if not e.get("is_special"):
            continue
        if e.get("reason") == "B":
            concept.append(e["term"])
        groups.setdefault(e.get("category") or "未分類", []).append(e["term"])
    n = sum(len(v) for v in groups.values())
    lines = ["---", f"domain: {domain}", f"term_count: {n}", "---", "",
             f"# {domain} 領域常駐詞表（CAG 前綴；收錄＝對照表 + Pass2 A/B/C）", ""]
    if concept:
        lines += ["## 概念術語（需固定譯法，最不能直譯錯）", "、".join(sorted(set(concept))), ""]
    for cat, ts in sorted(groups.items()):
        lines.append(f"## {cat}")
        lines.append("、".join(sorted(set(ts))))
        lines.append("")
    with open(os.path.join(_CAG, _safe(domain) + ".md"), "w", encoding="utf-8") as f:
        f.write("\n".join(lines) + "\n")
    return {"cag_terms": n, "file": _safe(domain) + ".md"}


def _acoustic_row(term: str, variants, is_special, tier) -> dict:
    return {"correct": term, "variants": [v for v in (variants or []) if v],
            "is_special": bool(is_special), "tier": tier or ""}


def _read_acoustic_rows() -> list[dict]:
    try:
        with open(_ACOUSTIC, encoding="utf-8") as f:
            return json.load(f).get("terms", [])
    except Exception:
        return []


def _save_acoustic(rows: list[dict]) -> dict:
    payload = {"count": len(rows),
               "variant_count": sum(len(r.get("variants") or []) for r in rows), "terms": rows}
    tmp = _ACOUSTIC + ".tmp"
    with open(tmp, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False)
    os.replace(tmp, _ACOUSTIC)
    return {"acoustic_terms": payload["count"], "acoustic_variants": payload["variant_count"]}


def _write_acoustic(terms: dict[str, dict]) -> dict:
    """全量寫聲學種子（§5.6 回灌）：知識側正解 term + ASR 變體 → 音韻索引消費。build 用。
    **只存原始字串、不烤 key**——音韻正規化單點走 IPA（`ipa.ipa_tokens`，phonetic_index 端算），
    符專案級全域正規化規範（不發散、改混淆集規則無舊 key 殘留）。"""
    rows = [_acoustic_row(e["term"], e.get("variants"), e.get("is_special"), e.get("tier"))
            for e in terms.values()]
    return _save_acoustic(rows)


def _acoustic_upsert(term: str, variants, is_special, tier) -> dict:
    """增量：新增／覆蓋 acoustic.json 裡 correct==term 那一筆，不重掃全 vault。"""
    rows = _read_acoustic_rows()
    row = _acoustic_row(term, variants, is_special, tier)
    for i, r in enumerate(rows):
        if r.get("correct") == term:
            rows[i] = row
            break
    else:
        rows.append(row)
    return _save_acoustic(rows)


def _acoustic_remove(term: str) -> dict:
    """增量：從 acoustic.json 移除 correct==term 那一筆（移出 CAG／刪／改名舊值用）。"""
    return _save_acoustic([r for r in _read_acoustic_rows() if r.get("correct") != term])


# ── CAG 變體人工 CRUD（前端直接改；改 vault md variants 行＋鎖定，刷 acoustic）───────────
def _terms_from_vault() -> dict[str, dict]:
    """從 vault terms md 重組 terms dict（給 _write_acoustic 刷新聲學種子）。"""
    from . import vault_graph
    out: dict[str, dict] = {}
    for n in vault_graph.graph(limit=1000000).get("nodes", []):
        if n.get("kind") == "term":
            out[n["id"]] = {"term": n["id"], "variants": n.get("variants", []),
                            "is_special": n.get("is_special", False), "tier": n.get("tier", "")}
    return out


def _invalidate_phonetic() -> None:
    try:
        from . import phonetic_index
        phonetic_index.invalidate()
    except Exception:
        pass


def refresh_acoustic() -> dict:
    """全量重刷 acoustic.json（從整個 vault；保留供全重建場景）。"""
    res = _write_acoustic(_terms_from_vault())
    _invalidate_phonetic()
    return res


def refresh_acoustic_one(term: str, removed: bool = False) -> dict:
    """增量刷單一 term 的聲學種子（人工 CRUD 用，取代全量 refresh_acoustic）：
    只讀該詞自己的 vault md → upsert acoustic.json 單筆，不重掃 198 詞。phonetic_index 下次查詢自動重建。
    removed=True（刪節點）或 md 已不存在 → 從聲學移除該筆。"""
    if removed:
        res = _acoustic_remove(term)
    else:
        path = _find_term_md(term)
        if not path:
            res = _acoustic_remove(term)
        else:
            meta, _ = _parse_term_note(path)
            res = _acoustic_upsert(meta.get("term") or term, meta.get("variants"),
                                   meta.get("is_special"), meta.get("tier"))
    _invalidate_phonetic()
    return res


def list_cag(q: str = "", only_special: bool = True) -> list[dict]:
    """列 CAG 詞 + 變體（前端變體管理用）。變體多的排前（雜訊好找）。"""
    from . import vault_graph
    out = []
    for n in vault_graph.graph(limit=1000000).get("nodes", []):
        if n.get("kind") != "term":
            continue
        if only_special and not n.get("is_special"):
            continue
        if q and q not in n["id"]:
            continue
        out.append({"term": n["id"], "variants": n.get("variants", []),
                    "category": n.get("category", ""), "tier": n.get("tier", ""),
                    "is_special": bool(n.get("is_special"))})
    out.sort(key=lambda r: (-len(r["variants"]), r["term"]))
    return out


def _find_term_md(term: str) -> str | None:
    """term → vault md 路徑（先試 _safe 檔名，再掃描比對 frontmatter term）。"""
    p = os.path.join(_VAULT, "terms", _safe(term) + ".md")
    if os.path.isfile(p):
        return p
    import glob

    from . import vault_graph
    for path in glob.glob(os.path.join(_VAULT, "terms", "*.md")):
        m = vault_graph._parse_note(path)
        if m and m.get("term") == term:
            return path
    return None


def set_variants(term: str, variants: list[str]) -> dict:
    """人工設定某 CAG 詞的變體（覆蓋）：改 vault md 的 variants 行 + locked:true（build 不沖），
    刷 acoustic.json + invalidate phonetic_index。回更新後狀態。找不到 md → 回 {error}。"""
    path = _find_term_md(term)
    if not path:
        return {"error": f"找不到 vault 詞：{term}"}
    cf = _flat(term)
    vs, seen = [], set()
    for v in variants or []:                          # 清理：去空白/去重/排除等於正解
        v = (v or "").strip()
        if v and _flat(v) != cf and _flat(v) not in seen:
            seen.add(_flat(v)); vs.append(v)
    with open(path, encoding="utf-8") as f:
        txt = f.read()
    txt = re.sub(r"^variants:.*$", "variants: [" + ", ".join(vs) + "]", txt, count=1, flags=re.M)
    if re.search(r"^locked:.*$", txt, flags=re.M):
        txt = re.sub(r"^locked:.*$", "locked: true", txt, count=1, flags=re.M)
    with open(path, "w", encoding="utf-8") as f:
        f.write(txt)
    ac = refresh_acoustic_one(term)       # 增量：只刷這一詞，不重掃全 vault
    return {"term": term, "variants": vs, "locked": True, **ac}


def set_term(term: str, is_special: bool | None = None, correct: str | None = None,
             delete: bool = False) -> dict:
    """人工校正 CAG 詞（鎖定、build 不還原）：
    - `is_special=False`＝標「不收」（移出 CAG/聲學索引，留心智圖）；`True`＝收回。
    - `correct`＝改正解（重命名 vault note + term 欄）。
    - `delete=True`＝刪整個心智圖節點（注意：來源仍在時 build 可能重抽；建議用『不收』持久排除）。"""
    path = _find_term_md(term)
    if not path:
        return {"error": f"找不到 vault 詞：{term}"}
    if delete:
        os.remove(path)
        return {"term": term, "deleted": True, **refresh_acoustic_one(term, removed=True)}
    with open(path, encoding="utf-8") as f:
        txt = f.read()
    if is_special is not None:
        txt = re.sub(r"^is_special:.*$", "is_special: " + str(bool(is_special)).lower(), txt, count=1, flags=re.M)
    new_path = path
    new_term = (correct or "").strip()
    if new_term and new_term != term:
        txt = re.sub(r"^term:.*$", "term: " + new_term, txt, count=1, flags=re.M)
        txt = re.sub(r"^# .*$", "# " + new_term, txt, count=1, flags=re.M)
        new_path = os.path.join(_VAULT, "terms", _safe(new_term) + ".md")
    if re.search(r"^locked:.*$", txt, flags=re.M):       # 人工校正一律鎖定，build 不還原
        txt = re.sub(r"^locked:.*$", "locked: true", txt, count=1, flags=re.M)
    with open(new_path, "w", encoding="utf-8") as f:
        f.write(txt)
    if new_path != path:
        os.remove(path)
        _acoustic_remove(term)            # 改名：先移除舊名聲學筆
    final = new_term or term
    return {"term": final, "is_special": is_special, "locked": True,
            **refresh_acoustic_one(final)}   # 增量：刷新／新增該詞單筆


def clear_outputs() -> None:
    """清空 vault(terms/domains)、cag、manifest（全重建前用；CAG 只由後端改）。"""
    import shutil
    for p in (os.path.join(_VAULT, "terms"), os.path.join(_VAULT, "domains")):
        if os.path.isdir(p):
            shutil.rmtree(p)
    if os.path.isdir(_CAG):
        for f in os.listdir(_CAG):
            if f.endswith(".md"):
                os.remove(os.path.join(_CAG, f))
    if os.path.exists(_MANIFEST):       # 增量快取也清，下次 build 全重抽
        os.remove(_MANIFEST)


# ── 後端可操控：列檔（含抽取狀態）/ 上傳 / 統計（前端儀表板用，§5.5）────────
_SUPPORTED = set(_CLEAN_EXT)            # 可被 build 收錄的副檔名（只乾淨來源；.srt/.txt 排除）


def list_sources() -> dict:
    """列 sources/ 全檔 + 抽取狀態（對 manifest 內容 hash 比對）+ 每檔抽取/排除決策。
    狀態：extracted(已抽、未變) / changed(內容改了，下次重抽) / new(新檔，下次抽) / skip(排除)。
    decision＝生效決策(extract|exclude)；override＝人工覆寫值或null；auto＝自動判斷；can_parse＝副檔有解析器。"""
    from collections import Counter
    manifest = _load_manifest()
    decisions = _load_decisions()
    out, cnt = [], Counter()
    for root, _dirs, files in os.walk(_SRC):
        # 列出 _excluded/ 等封存目錄（標排除讓前端看得到），只跳隱藏與 package(檢視產物)。build 仍不抽封存檔。
        _dirs[:] = [d for d in _dirs if not d.startswith(".") and d != "package"]
        rel_root = os.path.relpath(root, _SRC)
        archived = rel_root != "." and rel_root.split(os.sep)[0].startswith("_")   # 在 _xxx/ 底下＝封存排除
        for fn in sorted(files):
            if fn == ".gitkeep":
                continue
            path = os.path.join(root, fn)
            rel = os.path.relpath(path, _SRC)
            ext = os.path.splitext(fn)[1].lower()
            try:
                h = _file_hash(path)
                size = os.path.getsize(path)
            except OSError:
                continue
            cached = manifest.get(rel)
            eff, override = _effective_decision(rel, fn, ext, decisions)
            auto = _auto_decision(fn, ext)
            can_parse = ext == ".xlsx" or ext in _TEXT_EXT     # 有解析器（否則即使指定抽取也無效，如 .doc）
            if archived:                              # _excluded/ 封存 → 一律排除（不受人工 decision 影響）
                eff = auto = "exclude"
                override = None
                status = "skip"
            elif eff == "exclude":
                status = "skip"                       # 排除（自動或人工）
            elif cached and cached.get("hash") == h:
                status = "extracted"                  # 已抽、內容未變 → 走快取
            elif cached:
                status = "changed"                    # 內容改了 → 下次 build 重抽
            else:
                status = "new"                        # 新檔 → 下次 build 才抽
            cnt[status] += 1
            note = "封存（_excluded／移回 sources 根目錄才會抽）" if archived else (cached or {}).get("note", "")
            out.append({"rel": rel, "name": fn, "ext": ext, "size": size,
                        "status": status, "kind": (cached or {}).get("kind"), "note": note,
                        "archived": archived,
                        "decision": eff, "override": override, "auto": auto, "can_parse": can_parse})
    return {"sources": sorted(out, key=lambda r: (r["status"] != "new",
            r["status"] != "changed", r["rel"])),
            "counts": dict(cnt), "manifest_count": len(manifest), "overrides": len(decisions)}


def save_upload(filename: str, data: bytes, subdir: str = "") -> dict:
    """前端上傳：把檔寫進 sources/（可選 sub dir＝講者/場次）。不自動 build（讓使用者按鈕觸發）。"""
    name = os.path.basename(filename or "").strip()
    if not name:
        raise ValueError("檔名為空")
    sub = "/".join(p for p in (subdir or "").split("/") if p and p not in (".", ".."))
    dest_dir = os.path.join(_SRC, sub) if sub else _SRC
    os.makedirs(dest_dir, exist_ok=True)
    path = os.path.join(dest_dir, name)
    with open(path, "wb") as f:
        f.write(data)
    rel = os.path.relpath(path, _SRC)
    ext = os.path.splitext(name)[1].lower()
    supported = ext in _SUPPORTED and not _DIRTY_NAME.search(name)
    return {"rel": rel, "name": name, "size": len(data), "supported": supported,
            "note": "已存入 sources/，下次 build 會抽取" if supported else "副檔不支援或為逐字/字幕/校正類，build 會略過"}


def stats() -> dict:
    """知識庫儀表板統計（讀 vault note frontmatter 彙整，無 LLM）：總數/類別/tier/來源/特殊詞。"""
    from collections import Counter
    cats, tiers, srcs = Counter(), Counter(), Counter()
    total = special = 0
    tdir = os.path.join(_VAULT, "terms")
    if os.path.isdir(tdir):
        for fn in os.listdir(tdir):
            if not fn.endswith(".md"):
                continue
            meta, _ = _parse_term_note(os.path.join(tdir, fn))
            total += 1
            cats[meta.get("category") or "未分類"] += 1
            tiers[meta.get("tier") or "—"] += 1
            srcs[(meta.get("source") or "—").split(":")[0]] += 1
            if meta.get("is_special"):
                special += 1
    ddir = os.path.join(_VAULT, "domains")
    domains = len([f for f in os.listdir(ddir) if f.endswith(".md")]) if os.path.isdir(ddir) else 0
    cag = []
    if os.path.isdir(_CAG):
        for f in sorted(os.listdir(_CAG)):
            if not f.endswith(".md"):
                continue
            try:
                with open(os.path.join(_CAG, f), encoding="utf-8") as fh:
                    txt = fh.read()
                m = re.search(r"term_count:\s*(\d+)", txt)
                cag.append({"file": f, "term_count": int(m.group(1)) if m else 0})
            except OSError:
                pass
    return {"terms": total, "special": special, "domains": domains,
            "by_category": dict(cats.most_common()), "by_tier": dict(tiers.most_common()),
            "by_source": dict(srcs.most_common()), "cag": cag}


_REL_LINE = re.compile(r"關係： \[\[([^\]]+)\]\](?:（([^）]*)）)?")
_GLOSS_LINE = re.compile(r"^\*(.+)\*$", re.M)
_CTX_LINE = re.compile(r"^> (.+)$", re.M)


def _parse_term_note(path: str) -> tuple[dict, list]:
    """讀回 vault 詞 note → (entry, [(rel_target, desc)])。供 relabel 免重抽重建。"""
    with open(path, encoding="utf-8") as f:
        txt = f.read()
    meta = {"term": "", "term_lang": "", "ipa": "", "category": "", "source": "",
            "tier": "", "confidence": 0.0, "variants": [], "gloss": "", "context": "",
            "reason": "", "is_special": False, "locked": False, "en": "", "ja": "",
            "names": {}, "names_src": {}}
    if txt.startswith("---"):
        end = txt.find("\n---", 3)
        fm, body = txt[3:end], txt[end + 4:]
    else:
        fm, body = "", txt
    for ln in fm.splitlines():
        if ":" not in ln:
            continue
        k, v = ln.split(":", 1)
        k, v = k.strip(), v.strip()
        if k == "confidence":
            try:
                meta[k] = float(v)
            except ValueError:
                pass
        elif k == "variants":
            meta[k] = [x.strip() for x in v.strip("[]").split(",") if x.strip()]
        elif k in ("names", "names_src"):       # {eng: ..., jpn: ...} → dict
            meta[k] = {p.split(":", 1)[0].strip(): p.split(":", 1)[1].strip()
                       for p in v.strip("{}").split(",")
                       if ":" in p and p.split(":", 1)[1].strip()}
        elif k in ("is_special", "locked"):
            meta[k] = v.lower() == "true"
        elif k in meta:
            meta[k] = v
    g = _GLOSS_LINE.search(body)
    if g:
        meta["gloss"] = g.group(1).strip()
    c = _CTX_LINE.search(body)
    if c:
        meta["context"] = c.group(1).strip()          # 修 bug #4：relabel 不再遺失 context 例句
    rels = [(t.strip(), (d or "").strip()) for t, d in _REL_LINE.findall(body)]
    return meta, rels


def relabel(domain: str = "speakin") -> dict:
    """**免重抽**類別正規化：讀現有 vault → 類別併回 8 類實體 → 重生 terms/domains/cag。"""
    import glob
    import shutil
    from collections import Counter
    terms: dict[str, dict] = {}
    raw_rels: list[dict] = []
    for path in glob.glob(os.path.join(_VAULT, "terms", "*.md")):
        meta, rels = _parse_term_note(path)
        if not meta["term"]:
            continue
        meta["category"] = _canon_category(meta["category"])
        terms[_flat(meta["term"])] = meta
        for tgt, desc in rels:
            raw_rels.append({"source": meta["term"], "target": tgt, "desc": desc})
    rel_map: dict[tuple, dict] = {}
    for r in raw_rels:
        ks, kt = _flat(r["source"]), _flat(r["target"])
        if ks in terms and kt in terms and ks != kt:
            rel_map.setdefault((ks, kt), r)
    dp = os.path.join(_VAULT, "domains")
    if os.path.isdir(dp):
        shutil.rmtree(dp)                               # 舊類別 domains 全清，依新類別重生
    vault = _write_vault(terms, list(rel_map.values()))
    cag = _write_cag(terms, domain)
    acoustic = _write_acoustic(terms)                   # 同步刷新聲學種子
    return {"terms": len(terms), "categories": dict(Counter(e["category"] for e in terms.values())),
            "vault": vault, "cag": cag, "acoustic": acoustic}


async def build(domain: str = "speakin", use_llm: bool = True, llm_cap: int = 40,
                gleanings: int = 1, incremental: bool = True) -> dict:
    """主入口：sources/ → vault + cag。增量：只重抽新/改動檔（內容 hash），其餘讀快取；
    人工 locked 的 note 語意欄位不覆寫（§11 可持續擴充）。"""
    from collections import Counter
    h = await _harvest_sources(use_llm=use_llm, llm_cap=llm_cap, gleanings=gleanings,
                               incremental=incremental)
    terms, rels = h["terms"], h["relationships"]
    tune = await _tune_categories(terms) if use_llm else {"tuned": 0, "recovered": 0}  # 心智圖→tune misfit→CAG
    vault = _write_vault(terms, rels)
    cag = _write_cag(terms, domain)
    acoustic = _write_acoustic(terms)         # §5.6 回灌：正解+變體 → §4.1 音韻索引（下次查詢自動重建）
    return {"terms": len(terms), "relationships": len(rels), "reused_files": h.get("reused", 0),
            "tiers": dict(Counter(e.get("tier") for e in terms.values())),
            "categories": dict(Counter(e.get("category") for e in terms.values())),
            "tune": tune, "vault": vault, "cag": cag, "acoustic": acoustic,
            "used": h["used"], "skipped": h["skipped"]}
